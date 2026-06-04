"""
UC Davis College of Engineering PPT template deck builder.
Matches ECS 273 course presentation style (e.g. FinSight reference deck).
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = PROJECT_ROOT / "College of Engineering PPT Template - Clean.pptx"
TEMPLATE_FALLBACK = PROJECT_ROOT / "docs" / "ClariFi_Video_Presentation.pptx"
FOOTER_TEXT = "ECS 273 Spring Quarter 2026"
FONT = "Arial"

# Layout indices in template
LAYOUT_COVER = 0
LAYOUT_DIVIDER = 1
LAYOUT_TYPICAL = 6
LAYOUT_PIC_RIGHT = 5


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _font(paragraph, size: int, *, bold: bool = False) -> None:
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb("#000000")


def delete_all_slides(prs: Presentation) -> None:
    for sldId in list(prs.slides._sldIdLst):
        prs.part.drop_rel(sldId.rId)
        prs.slides._sldIdLst.remove(sldId)


def set_placeholder(slide, idx: int, text: str) -> None:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return


def add_course_footer(slide) -> None:
    box = slide.shapes.add_textbox(Inches(8.55), Inches(7.02), Inches(4.45), Inches(0.38))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = FOOTER_TEXT
    p.alignment = PP_ALIGN.RIGHT
    _font(p, 11)


def _fill_bullets(text_frame, lines: list[str], *, size: int = 20) -> None:
    text_frame.clear()
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line.strip()
        p.level = 0
        _font(p, size)
        p.space_after = Pt(10)


def content_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph.text_frame
        if ph.placeholder_format.idx == 2:
            return ph.text_frame
    return None


class EngineerDeck:
    def __init__(self, template_path: Path | None = None) -> None:
        path = template_path or TEMPLATE
        if not path.exists() and TEMPLATE_FALLBACK.exists():
            path = TEMPLATE_FALLBACK
        if not path.exists():
            raise FileNotFoundError(f"Template not found: {path}")
        self.template_path = path
        self.prs: Presentation | None = None

    def create(self, output_path: Path) -> None:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.output_path = output_path
        if output_path.resolve() != self.template_path.resolve():
            shutil.copy(self.template_path, output_path)
            self.prs = Presentation(str(output_path))
        else:
            self.prs = Presentation(str(self.template_path))
        delete_all_slides(self.prs)

    def save(self, output_path: Path | None = None) -> None:
        if self.prs is None:
            raise RuntimeError("Call create() first")
        self.prs.save(str(output_path or self.output_path))

    def add_cover(self, title: str, subtitle: str, notes: str = "") -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_COVER])
        set_placeholder(slide, 0, title)
        set_placeholder(slide, 10, subtitle)
        add_course_footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_section_divider(self, title: str, notes: str = "") -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_DIVIDER])
        set_placeholder(slide, 0, title)
        add_course_footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_bullet_slide(self, title: str, bullets: list[str], notes: str = "") -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_TYPICAL])
        set_placeholder(slide, 0, title)
        tf = content_placeholder(slide)
        if tf:
            _fill_bullets(tf, bullets[:8])
        add_course_footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_two_section_slide(
        self,
        title: str,
        section1_title: str,
        section1_bullets: list[str],
        section2_title: str,
        section2_bullets: list[str],
        notes: str = "",
    ) -> None:
        """Limitations & Future Scope style."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_TYPICAL])
        set_placeholder(slide, 0, title)
        lines = [section1_title, *section1_bullets, "", section2_title, *section2_bullets]
        tf = content_placeholder(slide)
        if tf:
            _fill_bullets(tf, lines, size=18)
        add_course_footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_image_right_slide(
        self,
        title: str,
        bullets: list[str],
        image_path: Path | None,
        notes: str = "",
    ) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_PIC_RIGHT])
        set_placeholder(slide, 0, title)
        tf = content_placeholder(slide)
        if tf:
            _fill_bullets(tf, bullets[:6], size=18)
        if image_path and image_path.exists():
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 10:
                    ph.insert_picture(str(image_path))
                    break
        add_course_footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_image_left_slide(
        self,
        title: str,
        bullets: list[str],
        image_path: Path | None,
        notes: str = "",
    ) -> None:
        """FinSight Visualisations layout: large image left, bullets right."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_TYPICAL])
        set_placeholder(slide, 0, title)
        if image_path and image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(0.45), Inches(1.25), width=Inches(6.35))
        box = slide.shapes.add_textbox(Inches(6.95), Inches(1.25), Inches(5.9), Inches(5.5))
        _fill_bullets(box.text_frame, bullets[:6], size=17)
        add_course_footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_full_image_slide(
        self,
        title: str,
        image_path: Path | None,
        caption: str = "",
        notes: str = "",
    ) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_TYPICAL])
        set_placeholder(slide, 0, title)
        if image_path and image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(0.55), Inches(1.15), width=Inches(12.2))
        if caption:
            cap = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(12.0), Inches(0.45))
            cap.text_frame.paragraphs[0].text = caption
            _font(cap.text_frame.paragraphs[0], 12)
        add_course_footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_work_distribution(
        self,
        title: str,
        members: list[str],
        tasks: list[tuple[str, list[bool]]],
        notes: str = "",
    ) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_TYPICAL])
        set_placeholder(slide, 0, title)
        nrows = len(tasks) + 1
        ncols = 1 + len(members)
        table = slide.shapes.add_table(
            nrows, ncols, Inches(0.55), Inches(1.35), Inches(12.2), Inches(5.2)
        ).table
        table.cell(0, 0).text = "Task Description"
        for c, name in enumerate(members, start=1):
            table.cell(0, c).text = name
        for r, (task, checks) in enumerate(tasks, start=1):
            table.cell(r, 0).text = task
            for c, on in enumerate(checks, start=1):
                table.cell(r, c).text = "✓" if on else ""
        for row in range(nrows):
            for col in range(ncols):
                for p in table.cell(row, col).text_frame.paragraphs:
                    _font(p, 14 if row == 0 else 13, bold=(row == 0))
        add_course_footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_thank_you(self, lines: list[str], notes: str = "") -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[LAYOUT_COVER])
        set_placeholder(slide, 0, "Thank You")
        set_placeholder(slide, 10, "\n".join(lines[:4]))
        add_course_footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
