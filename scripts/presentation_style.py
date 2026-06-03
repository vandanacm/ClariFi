"""
ECS273 conference-deck styling (presentation-styles.pdf).

- Light slides + dark text (content)
- ≤3 colors: teal accent, ink text, light background
- Sans serif (Helvetica / Arial)
- Meaningful headlines, ≤6 lines, prefer ≤3 bullets
- Slide numbers + outline (conference) + contact on last slide
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Three-color palette (+ white background)
TEAL = "#007f7a"
INK = "#1a2332"
MUTED = "#64748b"
BG = "#f8fafb"
WHITE = "#ffffff"

FONT = "Helvetica"
FONT_FALLBACK = "Arial"
TITLE_PT = 32
BODY_PT = 24
SUBTITLE_PT = 20
CAPTION_PT = 14
FOOTER_PT = 11
MAX_BODY_LINES = 6
PREFERRED_BULLETS = 3

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_font(run, size: int, *, bold: bool = False, color: str = INK) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def trim_bullets(bullets: list[str], max_lines: int = MAX_BODY_LINES) -> list[str]:
    out = [b.strip() for b in bullets if b and b.strip()]
    return out[:max_lines]


class ConferenceDeck:
    """Build conference-style slides (Liu / ECS273)."""

    def __init__(self, logo_path: Path | None = None) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.logo_path = logo_path
        self._slide_num = 0

    def _blank(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(BG)
        self._slide_num += 1
        return slide

    def _footer(self, slide, extra: str = "") -> None:
        text = f"{self._slide_num}" + (f"  ·  {extra}" if extra else "")
        box = slide.shapes.add_textbox(Inches(12.1), Inches(7.05), Inches(1.0), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        p.text = text
        set_font(p, FOOTER_PT, color=MUTED)

    def add_title_slide(
        self,
        title: str,
        subtitle: str,
        team_line: str = "",
        notes: str = "",
        *,
        dark: bool = True,
    ) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._slide_num += 1
        if dark:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb(INK)
            title_color, sub_color = WHITE, TEAL
        else:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb(BG)
            title_color, sub_color = INK, TEAL

        tf = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.5), Inches(3.0)).text_frame
        p = tf.paragraphs[0]
        p.text = title
        set_font(p, 44, bold=True, color=title_color)
        p2 = tf.add_paragraph()
        p2.text = subtitle
        set_font(p2, SUBTITLE_PT, color=sub_color)
        if team_line:
            p3 = tf.add_paragraph()
            p3.text = team_line
            set_font(p3, 18, color=MUTED if not dark else WHITE)

        if self.logo_path and self.logo_path.exists():
            slide.shapes.add_picture(str(self.logo_path), Inches(11.0), Inches(0.4), height=Inches(0.85))

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_outline_slide(self, sections: list[str], notes: str = "") -> None:
        slide = self._blank()
        self._headline(slide, "Outline")
        bullets = trim_bullets(sections, max_lines=6)
        self._bullets(slide, bullets, top=Inches(1.35))
        self._footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_content_slide(
        self,
        headline: str,
        bullets: list[str],
        *,
        image_path: Path | None = None,
        caption: str = "",
        notes: str = "",
        image_only: bool = False,
    ) -> None:
        slide = self._blank()
        self._headline(slide, headline)
        bullets = trim_bullets(bullets)

        if image_only and image_path and image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(0.65), Inches(1.05), width=Inches(12.05))
            if caption:
                cap = slide.shapes.add_textbox(Inches(0.75), Inches(6.75), Inches(11.5), Inches(0.45))
                cap.text_frame.paragraphs[0].text = caption
                set_font(cap.text_frame.paragraphs[0], CAPTION_PT, color=MUTED)
        else:
            left_w = Inches(5.6) if image_path and image_path.exists() else Inches(12.0)
            self._bullets(slide, bullets, top=Inches(1.35), width=left_w)
            if image_path and image_path.exists():
                slide.shapes.add_picture(str(image_path), Inches(6.35), Inches(1.2), width=Inches(6.5))

        self._footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_contact_slide(
        self,
        project: str,
        team_lines: list[str],
        contact: str = "your.email@ucdavis.edu",
        notes: str = "",
    ) -> None:
        slide = self._blank()
        self._headline(slide, "Thank you — questions?")
        lines = [project, *team_lines[:3], contact]
        self._bullets(slide, lines, top=Inches(1.5), width=Inches(11.0))
        self._footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_table_slide(
        self,
        headline: str,
        headers: list[str],
        rows: list[list[str]],
        notes: str = "",
    ) -> None:
        slide = self._blank()
        self._headline(slide, headline)
        nrows, ncols = len(rows) + 1, len(headers)
        table = slide.shapes.add_table(
            nrows, ncols, Inches(1.0), Inches(1.45), Inches(11.3), Inches(4.8)
        ).table
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            set_font(cell.text_frame.paragraphs[0], BODY_PT - 2, bold=True)
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.text = val
                set_font(cell.text_frame.paragraphs[0], BODY_PT - 4)
        self._footer(slide)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def save(self, path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))

    def _headline(self, slide, text: str) -> None:
        tb = slide.shapes.add_textbox(Inches(0.65), Inches(0.4), Inches(12.0), Inches(0.75))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        set_font(p, TITLE_PT, bold=True, color=TEAL)

    def _bullets(
        self,
        slide,
        bullets: list[str],
        *,
        top=Inches(1.35),
        width=Inches(12.0),
    ) -> None:
        body = slide.shapes.add_textbox(Inches(0.65), top, width, Inches(5.6))
        tf = body.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            set_font(p, BODY_PT)
            p.space_after = Pt(12)
